from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Iterable, List, Optional

try:
    import nbformat as nbf
except Exception:  # pragma: no cover
    nbf = None  # Will raise a friendly error when used


LECTURE_TITLES = [
    "01 - What is MLOps + Local Stack Setup",
    "02 - Reproducibility and Configuration",
    "03 - Data Management and Versioning",
    "04 - Data Quality and Validation",
    "05 - Experiment Tracking",
    "06 - Project Structure and Testing for ML",
    "07 - Orchestration and Pipelines",
    "08 - CI for ML with GitHub Actions",
    "09 - Model Packaging and Registry",
    "10 - Serving: Batch and Real-time",
    "11 - Monitoring and Observability",
    "12 - Deployment Strategies",
    "13 - Security, Governance, and Responsible AI",
    "14 - Capstone, SLAs, Incident Response, and Postmortems",
]

DEFAULT_PIP_PACKAGES = [
    "pip",
    "setuptools",
    "wheel",
    "numpy",
    "pandas",
    "matplotlib",
    "seaborn",
    "scikit-learn",
    "mlflow",
    "dvc",
    "prefect",
    "fastapi",
    "uvicorn",
    "pydantic",
    "hydra-core",
    "great-expectations",
    "evidently",
]


def _ensure_dir(path: Path) -> None:
    """Create a directory if it does not exist.

    Parameters
    ----------
    path : Path
        Directory path to ensure exists.
    """
    path.mkdir(parents=True, exist_ok=True)


def _lecture_objectives(idx: int) -> List[str]:
    """Return short, actionable objectives for a given lecture.

    Parameters
    ----------
    idx : int
        Lecture index starting at 1.

    Returns
    -------
    list of str
        Objectives bullet points.
    """
    mapping = {
        1: [
            "Understand the ML system lifecycle and the role of MLOps.",
            "Install and verify the local development stack.",
            "Create a clean repository skeleton.",
        ],
        2: [
            "Pin environments and ensure deterministic runs.",
            "Manage configs with Hydra and validate with Pydantic.",
        ],
        3: [
            "Track datasets and pipelines with DVC.",
            "Implement raw -> processed data steps and splits.",
        ],
        4: [
            "Define data contracts and expectations.",
            "Add data tests and CI checks.",
        ],
        5: [
            "Track experiments with MLflow.",
            "Log params, metrics, artifacts; compare runs.",
        ],
        6: [
            "Structure ML projects as installable packages.",
            "Write fast tests with pytest and pre-commit hooks.",
        ],
        7: [
            "Build parameterized pipelines with Prefect.",
            "Add caching, retries, and artifact handoffs.",
        ],
        8: [
            "Design CI for data-dependent ML workflows.",
            "Build and cache environments; run tests and lint.",
        ],
        9: [
            "Package and register models with MLflow Models.",
            "Manage model stages and promotions.",
        ],
        10: [
            "Serve models via FastAPI for batch and real-time.",
            "Containerize and expose health checks.",
        ],
        11: [
            "Monitor drift and performance with Evidently.",
            "Add logging/metrics for observability.",
        ],
        12: [
            "Understand deployment strategies and rollbacks.",
            "Operate staging vs. prod with docker-compose.",
        ],
        13: [
            "Handle secrets, PII, and supply-chain risks.",
            "Produce model cards and lineage.",
        ],
        14: [
            "Define SLOs/SLAs and handle incidents.",
            "Deliver capstone demo and postmortem.",
        ],
    }
    return mapping.get(idx, ["See course README for details."])


def _pip_install_cell(packages: Iterable[str]) -> str:
    """Build a portable pip install cell.

    Parameters
    ----------
    packages : Iterable[str]
        Package names to install with pip.

    Returns
    -------
    str
        Python code for a Jupyter cell that installs packages.
    """
    pkgs_json = json.dumps(list(packages))
    return f"""# Install runtime dependencies (idempotent). Safe to re-run.
import sys, subprocess, json
pkgs = json.loads('''{pkgs_json}''')
# Upgrade pip tooling first for reliable wheels
subprocess.check_call([sys.executable, "-m", "pip", "install", "-q", "--upgrade", "pip", "setuptools", "wheel"])
# Install core packages
subprocess.check_call([sys.executable, "-m", "pip", "install", "-q"] + pkgs)
print("Installed packages:", ", ".join(pkgs))
"""


def _notebook_cells(title: str, idx: int, packages: Iterable[str]) -> List:
    """Construct default cells for a lecture notebook.

    Parameters
    ----------
    title : str
        Notebook title text.
    idx : int
        Lecture index starting at 1.
    packages : Iterable[str]
        Pip packages to install in the first code cell.

    Returns
    -------
    list
        A list of nbformat cells.
    """
    if nbf is None:  # pragma: no cover
        raise RuntimeError("nbformat is required to create notebooks. Please install nbformat.")

    md_title = nbf.v4.new_markdown_cell(f"# {title}\n\nThis notebook contains hands-on materials for lecture {idx:02d}.")
    md_objectives = nbf.v4.new_markdown_cell(
        "## Objectives\n\n" + "\n".join([f"- {o}" for o in _lecture_objectives(idx)])
    )
    cell_pip = nbf.v4.new_code_cell(_pip_install_cell(packages))
    cell_imports = nbf.v4.new_code_cell(
        "import numpy as np\nimport pandas as pd\nimport matplotlib.pyplot as plt\n"
        "from pathlib import Path\n\nprint('Libraries imported.')"
    )
    cell_template = nbf.v4.new_markdown_cell(
        "## Your tasks\n\n"
        "- Implement the steps for this lecture below.\n"
        "- Keep cells small and testable.\n"
        "- Save any generated figures in variables for reuse (e.g., `fig`)."
    )
    cell_example = nbf.v4.new_code_cell(
        "## Example: quick sanity plot\n"
        "import numpy as np\nimport matplotlib.pyplot as plt\n\n"
        "x = np.linspace(0, 2*np.pi, 200)\n"
        "y = np.sin(x)\n"
        "fig, ax = plt.subplots()\nax.plot(x, y)\nax.set(title='Sanity Plot: Sine Wave')\n"
        "plt.show()\n"
        "fig  # Keep reference for later reuse"
    )
    return [md_title, md_objectives, cell_pip, cell_imports, cell_template, cell_example]


def create_notebook(path: Path, title: str, idx: int, packages: Optional[Iterable[str]] = None) -> None:
    """Create a Jupyter notebook with a first cell that installs requirements.

    Parameters
    ----------
    path : Path
        Destination .ipynb path.
    title : str
        Notebook title for the lecture.
    idx : int
        Lecture index starting at 1.
    packages : iterable of str, optional
        Additional packages to install. Defaults to a standard MLOps stack.

    Notes
    -----
    - The first code cell installs a minimal set of MLOps packages.
    - The notebook includes an example plot whose figure is stored in a variable `fig`.
    """
    if nbf is None:  # pragma: no cover
        raise RuntimeError("nbformat is required to create notebooks. Please install nbformat.")

    pkgs = list(packages) if packages is not None else DEFAULT_PIP_PACKAGES
    nb = nbf.v4.new_notebook()
    nb["cells"] = _notebook_cells(title=title, idx=idx, packages=pkgs)
    _ensure_dir(path.parent)
    with path.open("w", encoding="utf-8") as f:
        nbf.write(nb, f)


def create_slide_deck(path: Path, title: str, bullets: Iterable[str]) -> None:
    """Create a slide deck for a lecture.

    Tries to generate a PowerPoint (.pptx) using python-pptx if available.
    Falls back to a Markdown slide outline if python-pptx is not installed.

    Parameters
    ----------
    path : Path
        Destination path. Should ideally end with .pptx. If python-pptx is not available,
        a .md file with the same stem is created.
    title : str
        Deck title.
    bullets : iterable of str
        Bullet points to include on the first slide.
    """
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore

        _ensure_dir(path.parent)
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        subtitle = slide.placeholders[1]
        subtitle.text = "Applied MLOps - Lecture Materials"

        content_slide_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(content_slide_layout)
        slide2.shapes.title.text = "Objectives"
        body = slide2.shapes.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(bullets):
            p = body.add_paragraph() if i else body.paragraphs[0]
            p.text = str(b)
            p.level = 0

        prs.save(path)
    except Exception:
        # Fallback to Markdown
        md_path = path.with_suffix(".md")
        _ensure_dir(md_path.parent)
        lines = [f"# {title}", "", "## Objectives", *[f"- {b}" for b in bullets]]
        md_path.write_text("\n".join(lines), encoding="utf-8")


def scaffold(base_dir: Path = Path(".")) -> None:
    """Create a minimal starter repo skeleton with notebooks and slides.

    Parameters
    ----------
    base_dir : Path, optional
        Base directory where the scaffold is created. Defaults to current directory.

    Creates
    -------
    - notebooks/lec_XX_title.ipynb for 14 lectures
    - slides/lec_XX_title.pptx (or .md if python-pptx is missing)
    - directories: configs/, data/, docker/, tests/, .github/workflows/
    - placeholder workflow file if none exists
    """
    # Core folders
    for d in ["notebooks", "slides", "configs", "data", "docker", "tests", ".github/workflows"]:
        _ensure_dir(base_dir / d)

    # Notebooks and slides
    for i, title in enumerate(LECTURE_TITLES, start=1):
        safe_title = title.split(" - ")[1].lower().replace(" ", "_").replace("/", "_")
        nb_path = base_dir / "notebooks" / f"lec_{i:02d}_{safe_title}.ipynb"
        create_notebook(nb_path, title=title, idx=i, packages=DEFAULT_PIP_PACKAGES)

        slide_path = base_dir / "slides" / f"lec_{i:02d}_{safe_title}.pptx"
        create_slide_deck(slide_path, title=title, bullets=_lecture_objectives(i))

    # Minimal GH Actions CI if missing
    ci_path = base_dir / ".github" / "workflows" / "ci.yml"
    if not ci_path.exists():
        ci_yaml = """name: ci
on:
  push:
  pull_request:
jobs:
  lint-test:
    runs-on: ubuntu-latest
    steps:
      - uses: actions/checkout@v4
      - uses: actions/setup-python@v5
        with:
          python-version: '3.10'
      - name: Install deps
        run: |
          python -m pip install --upgrade pip
          pip install pytest black isort flake8 mypy
      - name: Lint
        run: |
          black --check .
          isort --check-only .
          flake8 .
      - name: Type check
        run: mypy --ignore-missing-imports src || true
      - name: Tests
        run: pytest -q || true
"""
        ci_path.write_text(ci_yaml, encoding="utf-8")


def parse_args(argv: Optional[List[str]] = None) -> argparse.Namespace:
    """Parse CLI arguments.

    Parameters
    ----------
    argv : list of str, optional
        Argument vector. If None, uses sys.argv.

    Returns
    -------
    argparse.Namespace
        Parsed arguments.
    """
    p = argparse.ArgumentParser(description="MLOps course scaffolder")
    sub = p.add_subparsers(dest="cmd", required=True)

    s_scaffold = sub.add_parser("scaffold", help="Create notebooks and slides for the 14-lecture MLOps course")
    s_scaffold.add_argument("--base-dir", type=Path, default=Path("."), help="Target directory (default: current dir)")

    return p.parse_args(argv)


def main(argv: Optional[List[str]] = None) -> None:
    """Entry point for the CLI.

    Parameters
    ----------
    argv : list of str, optional
        CLI arguments. If None, uses sys.argv.
    """
    args = parse_args(argv)
    if args.cmd == "scaffold":
        scaffold(args.base_dir)


if __name__ == "__main__":  # pragma: no cover
    main()
