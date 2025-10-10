"""MLOps course starter generator.

This module provides a tiny utility to scaffold a semester-long, applied MLOps
starter repository with:
- notebooks/ with one Jupyter Notebook per lecture (incl. a pip-install cell)
- slides/ with a PowerPoint per lecture (or a Markdown fallback if python-pptx is missing)

Run create_starter_repo() from your project root to generate files.

"""

from __future__ import annotations

import json
import os
from datetime import date
from typing import Dict, Iterable, List, Sequence, Tuple

__all__ = ["create_starter_repo", "COURSE_OUTLINE"]
__version__ = "0.1.0"


COURSE_OUTLINE: List[Tuple[str, str]] = [
    ("01_intro", "What is MLOps + local stack setup"),
    ("02_reproducibility", "Reproducibility and configuration"),
    ("03_data_versioning", "Data management and versioning"),
    ("04_data_quality", "Data quality and validation"),
    ("05_tracking", "Experiment tracking"),
    ("06_project_structure", "Project structure and testing for ML"),
    ("07_orchestration", "Orchestration and pipelines"),
    ("08_ci", "CI for ML with GitHub Actions"),
    ("09_registry", "Model packaging and registry"),
    ("10_serving", "Serving: batch and real-time"),
    ("11_monitoring", "Monitoring and observability"),
    ("12_deployment", "Deployment strategies"),
    ("13_governance", "Security, governance, and responsible AI"),
    ("14_capstone", "Capstone, SLAs, incident response, and postmortems"),
]

_BASE_PACKAGES: Sequence[str] = (
    "numpy",
    "pandas",
    "scikit-learn",
    "mlflow",
    "dvc",
    "prefect>=2",
    "fastapi",
    "uvicorn[standard]",
    "great-expectations",
    "evidently",
    "pydantic>=1.10,<3",
    "hydra-core",
    "python-dotenv",
    "pytest",
    "black",
    "isort",
    "flake8",
    "mypy",
)


def _ensure_dir(path: str) -> None:
    os.makedirs(path, exist_ok=True)


def _write_text(path: str, text: str) -> None:
    _ensure_dir(os.path.dirname(path))
    with open(path, "w", encoding="utf-8") as f:
        f.write(text)


def _write_bytes(path: str, data: bytes) -> None:
    _ensure_dir(os.path.dirname(path))
    with open(path, "wb") as f:
        f.write(data)


def _make_notebook(title: str, packages: Sequence[str]) -> Dict:
    """Create a minimal Jupyter notebook JSON with a pip-install cell.

    Parameters
    ----------
    title : str
        Title shown in the first markdown cell.
    packages : Sequence[str]
        Packages that will be pip-installed in the environment cell.

    Returns
    -------
    dict
        Notebook JSON structure (nbformat v4).
    """
    install_code = (
        "import sys, subprocess\n"
        "pkgs = [\n"
        + "".join(f"    '{p}',\n" for p in packages)
        + "]\n"
        "subprocess.check_call([sys.executable, '-m', 'pip', 'install', '-q'] + pkgs)\n"
        "print('Installed:', ', '.join(pkgs))\n"
    )

    info_code = (
        "import sys, platform\n"
        "import numpy as np, pandas as pd\n"
        "print('Python', sys.version)\n"
        "print('Platform', platform.platform())\n"
        "print('NumPy', np.__version__)\n"
        "print('Pandas', pd.__version__)\n"
    )

    body_md = (
        "## Objectives\n"
        "- Follow the checklist in the course README for this lecture.\n"
        "- Replace this cell with your hands-on work.\n"
    )

    nb = {
        "cells": [
            {"cell_type": "markdown", "metadata": {}, "source": f"# {title}\nGenerated: {date.today().isoformat()}"},
            {"cell_type": "code", "metadata": {"tags": ["setup"]}, "execution_count": None, "outputs": [], "source": install_code},
            {"cell_type": "code", "metadata": {}, "execution_count": None, "outputs": [], "source": info_code},
            {"cell_type": "markdown", "metadata": {}, "source": body_md},
            {"cell_type": "code", "metadata": {}, "execution_count": None, "outputs": [], "source": "# Your code here\n"},
        ],
        "metadata": {
            "kernelspec": {"display_name": "Python 3", "language": "python", "name": "python3"},
            "language_info": {"name": "python"},
        },
        "nbformat": 4,
        "nbformat_minor": 5,
    }
    return nb


def _save_notebook(path: str, nb: Dict) -> None:
    _ensure_dir(os.path.dirname(path))
    with open(path, "w", encoding="utf-8") as f:
        json.dump(nb, f, ensure_ascii=False, indent=2)


def _make_pptx(title: str) -> bytes:
    """Create a minimal PowerPoint file.

    Tries to use python-pptx; falls back to a simple placeholder markdown if not installed.

    Parameters
    ----------
    title : str
        Slide deck title.

    Returns
    -------
    bytes
        PPTX file bytes; if python-pptx is unavailable, returns empty bytes.
    """
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"{date.today().isoformat()} • MLOps"
        # Basic second slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Objectives"
        slide2.placeholders[1].text = "- Overview\n- Hands-on tasks\n- Deliverables"
        from io import BytesIO

        bio = BytesIO()
        prs.save(bio)
        return bio.getvalue()
    except Exception:
        return b""


def create_starter_repo(base_dir: str = ".", make_slides: bool = True) -> Dict[str, List[str]]:
    """Generate a minimal MLOps course starter in the given directory.

    This creates:
    - notebooks/lecture_XX_title.ipynb for each lecture with a pip install cell
    - slides/lecture_XX_title.pptx if python-pptx is available, otherwise a .md placeholder

    Parameters
    ----------
    base_dir : str, default="."
        Directory where the starter structure will be generated.
    make_slides : bool, default=True
        Whether to generate slide decks.

    Returns
    -------
    dict
        Mapping with lists of created files under keys: notebooks, slides, misc.

    Examples
    --------
    >>> from mlops import create_starter_repo
    >>> results = create_starter_repo(".")
    >>> results["notebooks"][:2]  # doctest: +SKIP
    ['notebooks/lecture_01_intro.ipynb', 'notebooks/lecture_02_reproducibility.ipynb']
    """
    created_nb: List[str] = []
    created_slides: List[str] = []
    created_misc: List[str] = []

    notebooks_dir = os.path.join(base_dir, "notebooks")
    slides_dir = os.path.join(base_dir, "slides")
    _ensure_dir(notebooks_dir)
    if make_slides:
        _ensure_dir(slides_dir)

    # Notebooks
    for idx, (slug, title) in enumerate(COURSE_OUTLINE, start=1):
        fname = f"lecture_{idx:02d}_{slug}.ipynb"
        path = os.path.join(notebooks_dir, fname)
        nb = _make_notebook(f"Lecture {idx}: {title}", _BASE_PACKAGES)
        _save_notebook(path, nb)
        created_nb.append(os.path.relpath(path, base_dir))

    # Slides
    if make_slides:
        for idx, (slug, title) in enumerate(COURSE_OUTLINE, start=1):
            pptx_name = f"lecture_{idx:02d}_{slug}.pptx"
            pptx_path = os.path.join(slides_dir, pptx_name)
            data = _make_pptx(f"Lecture {idx}: {title}")
            if data:
                _write_bytes(pptx_path, data)
                created_slides.append(os.path.relpath(pptx_path, base_dir))
            else:
                md_path = os.path.join(slides_dir, f"lecture_{idx:02d}_{slug}.md")
                _write_text(
                    md_path,
                    f"# Lecture {idx}: {title}\n\nInstall python-pptx and re-run create_starter_repo to generate PPTX.\n",
                )
                created_slides.append(os.path.relpath(md_path, base_dir))

    # Minimal README snippet (non-destructive if exists)
    readme_path = os.path.join(base_dir, "README.md")
    if not os.path.exists(readme_path):
        _write_text(
            readme_path,
            "# MLOps Course Starter\n\n- See notebooks/ for per-lecture hands-on notebooks.\n- See slides/ for lecture slides.\n",
        )
        created_misc.append(os.path.relpath(readme_path, base_dir))

    return {"notebooks": created_nb, "slides": created_slides, "misc": created_misc}
